import demistomock as demisto  # noqa: F401
from CommonServerPython import *  # noqa: F401
import openpyxl
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
import pandas as pd
from docx.opc.constants import RELATIONSHIP_TYPE as RT


def extract_hyperlinks_from_xlsx(file_path: str) -> Set:
    with zipfile.ZipFile(file_path, "r") as zf:
        xmls = [zf.read(fn) for fn in zf.infolist()
                if fn.filename.startswith("xl/drawings/_rels/")]

    urls = set()

    for xml_data in xmls:
        df = pd.read_xml(xml_data)

        if "TargetMode" in df.columns:
            filtered_df = df.loc[df["TargetMode"].eq("External"), "Target"]
            urls |= set(filtered_df)

    wb = openpyxl.load_workbook(file_path)
    for sheet in wb:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.hyperlink:
                    urls.add(cell.hyperlink.target)

    return urls


def extract_hyperlinks_from_docx(file_path: str) -> Set:
    doc = Document(file_path)
    links = set()
    for rel in doc.part.rels.values():
        if rel.reltype == RT.HYPERLINK and rel.is_external:
            links.add(rel._target)

    return links


def extract_hyperlinks_from_pptx(file_path: str) -> Set:
    prs = Presentation(file_path)
    links = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            links.add(run.hyperlink.address)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:    # pylint: disable=E1101
                group_shape = shape
                for s in group_shape.shapes:
                    if s.click_action and s.click_action.hyperlink.address:
                        links.add(s.click_action.hyperlink.address)
            elif shape.click_action and shape.click_action.hyperlink.address:
                links.add(shape.click_action.hyperlink.address)

    return links


def extract_hyperlink_by_file_type(file_name: str, file_path: str) -> CommandResults:
    if file_name.lower().endswith('.xlsx'):
        result = extract_hyperlinks_from_xlsx(file_path)
    elif file_name.lower().endswith('.docx'):
        result = extract_hyperlinks_from_docx(file_path)
    elif file_name.lower().endswith('.pptx'):
        result = extract_hyperlinks_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file type. Supported types are: 'xlsx, docx, pptx'")
    if result:
        urls_str = "\n".join(result)
        hr = f'### Extracted Hyperlinks\n\n{urls_str}'
    else:
        hr = '**No hyperlinks.**'

    output = [{'URL': url, 'FileName': file_name} for url in result]
    return CommandResults(
        outputs=output,
        outputs_prefix='ExtractedHyperLink',
        outputs_key_field=['URL', 'FileName'],
        readable_output=hr,
        raw_response=list(result)
    )


def main():  # pragma: no cover
    try:
        entry_id = demisto.args().get("entry_id")
        file_result = demisto.getFilePath(entry_id)
        if not file_result:
            raise ValueError(f"Couldn't find entry id: {entry_id}")
        file_name = file_result.get('name')
        file_path = file_result.get('path')
        os.rename(f'./{file_path}', file_name)
        return_results(extract_hyperlink_by_file_type(file_name=file_name, file_path=os.path.realpath(file_name)))
    except Exception as ex:
        return_error(f'Failed to execute ExtractHyperlinksFromOfficeFiles. Error: {str(ex)}')


if __name__ in ('__main__', '__builtin__', 'builtins'):
    main()
